from pptx import Presentation

prs = Presentation("DE_Proyecto_Avance.pptx")
print(f"Total slides: {len(prs.slides)}")

for i, slide in enumerate(prs.slides):
    print(f"\n--- Slide {i+1} ---")
    if slide.shapes.title:
        print(f"Title: {slide.shapes.title.text}")
    else:
        print("No Title Shape")
        
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if len(run.text.strip()) > 0:
                        print(f"  [Text]: {run.text}")
