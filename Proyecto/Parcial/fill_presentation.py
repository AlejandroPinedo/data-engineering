from pptx import Presentation
import os

TEMPLATE_PATH = "DE_Proyecto_Avance.pptx"
OUTPUT_PATH = "Consignacion_Mina_Avance_Proyecto.pptx"

if not os.path.exists(TEMPLATE_PATH):
    print(f"[ERROR] No se encontro la plantilla base: {TEMPLATE_PATH}")
    exit(1)

prs = Presentation(TEMPLATE_PATH)

# Tabla de reemplazos de texto clave de placeholders a datos reales
reemplazos = {
    # Portada / Integrantes
    "Integrantes :": "Integrantes :\n- [Integrante 1]\n- [Integrante 2]\n- [Integrante 3]\n- [Integrante 4]\n- [Integrante 5]",
    
    # Empresa y Rubro
    "NOMBRE DE LA EMPRESA": "Compañía Minera Las Bambas S.A.",
    "RUBRO o INDUSTRIA": "Minería - Logística y Cadena de Suministro",
    
    # Datasets
    "Ejemplos - Datasets por rubro": "Dataset de Almacén: Stock de Materiales y Repuestos en Consignación\n(CSV diario exportado por SAP con variables de id_material, descripción, proveedor, cantidades, mínimos de seguridad y ubicación)",
    
    # Introducción y Problema
    "Contexto del dominio.": "Contexto: Control de inventarios en consignación en almacenes mineros en los Andes peruanos, con alta dependencia de repuestos críticos de proveedores como Ferreyros y Komatsu.",
    "¿Qué problema quieren resolver?": "Problema: Desabastecimiento de repuestos críticos (stockouts) debido a reportes manuales semanales por correo, deteniendo operaciones mineras con pérdidas de $10K-$50K USD por hora.",
    "¿Por qué importa automatizar esta carga?": "Automatización: Proveer visibilidad en tiempo real a los proveedores sobre el stock físico actual en mina para gatillar reposiciones automáticas y eliminar el desabastecimiento.",
    
    # Objetivo
    "Qué se busca lograr en esta primera entrega.": "Objetivo: Ingestar de forma automática y robusta los reportes diarios de SAP a Azure Data Lake (contenedor raw) orquestado mediante Airflow.",
    "Cuál será el alcance (pipeline hasta Raw zone, sin transformación).": "Alcance: Ingesta segura e inmutable desde el servidor local simulado hasta la zona Raw de almacenamiento seguro en nube, sentando las bases del procesamiento Medallion.",
    
    # Arquitectura
    "Local -> Airflow -> Datalake": "SAP Local (Mina) -> Apache Airflow (Orquestador) -> Azure Data Lake (Contenedor Raw)",
    "Fuentes de datos": "Fuentes de datos: CSVs diarios exportados por SAP con stock físico, niveles mínimos de seguridad y costos unitarios.",
    "Orquestador": "Orquestador: Apache Airflow (local/dockerizado) que controla la periodicidad diaria, reintentos y alertas.",
    "Zona de almacenamiento": "Zona de almacenamiento: Contenedor 'raw' en Azure Data Lake Storage (ADLS) Gen2.",
    "Por qué se eligieron esos componentes?": "Justificación: ADLS Gen2 ofrece almacenamiento en la nube resiliente e inmutable, y Airflow asegura control del flujo diario, reintentos automatizados ante fallas de red y alertamiento.",
    
    # DAG
    "¿Qué tareas hay en el DAG?": "DAG: telco_mina_consignacion_ingestion_dag\nTareas: esperar_csv_sap_consignacion (FileSensor) -> subir_a_azure_y_archivar_consignacion (PythonOperator con WasbHook y archivado local).",
    "¿Con qué frecuencia corre?": "Frecuencia: Ejecución diaria (@daily) tras el cierre de operaciones de almacén a la medianoche.",
    "¿Qué errores se contemplan?": "Manejo de errores: Sensor configurado con timeout de 5 min. Reintentos automáticos del pipeline si falla la conexión con Azure. Transaccionalidad de archivos (solo se mueven a archive/ tras subida exitosa).",
    
    # Git
    "Enlace al repositorio": "Repositorio: GitHub (dsai-de-airflow)",
    "Estrategia Gitflow": "Estrategia: Gitflow, aislando los desarrollos nuevos en ramas feature/",
    "Ramas utilizadas": "Ramas utilizadas: main (producción), develop (integración), feature/consignment-inventory-ingestion (desarrollo)",
    "Organización básica del código": "Organización: Carpetas estructuradas (dags/, datos_simulados/) y README.md que documenta la instalación y despliegue del proyecto.",
    
    # Validación
    "¿El DAG corre correctamente?": "Validación: Ejecución local exitosa con 100% de éxito en la subida y archivado.",
    "¿Qué outputs se generan en la zona Raw?": "Outputs en raw: Archivo inmutable en raw/inventario_consignacion/sap_consignacion_YYYYMMDD_HHMMSS.csv.",
    "(Opcional: capturas de la interfaz de Airflow y carpeta destino)": "[Espacio reservado para capturas de pantalla de la UI de Airflow y del Portal de Azure Storage]",
    
    # Desafíos
    "¿Qué fue lo más difícil?": "Desafío: Mapear correctamente los volúmenes del sistema de archivos local con el contenedor Docker de Airflow para que el FileSensor detecte el archivo de forma confiable.",
    "¿Qué aprendieron al resolverlo?": "Aprendizaje: Importancia de la transaccionalidad al subir y mover archivos locales en una sola tarea atómica para evitar duplicaciones.",
    
    # Siguiente fase
    "Cómo conectarán esto con el Data Product futuro?": "Data Product: Tablero de control de inventario en consignación compartido con proveedores y sistema de alertas automáticas de reposición.",
    "¿Qué tipo de tablas podría surgir en Silver/Gold?": "Bronze: Histórico consolidado 1:1.\nSilver: Deduplicación temporal, tipado y cuarentena de mediciones corruptas (stock negativo).\nGold: Datamart de reposición rápida (stock < seguridad)."
}

# Recorrer slides y shapes para hacer reemplazo de texto
reemplazados_count = 0
for i, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    orig_text = run.text
                    for placeholder, reemplazo in reemplazos.items():
                        if placeholder in orig_text:
                            run.text = orig_text.replace(placeholder, reemplazo)
                            reemplazados_count += 1
                            print(f"[REEMPLAZO] Slide {i+1}: '{placeholder}' -> '{reemplazo[:40]}...'")

# Guardar la nueva presentación
prs.save(OUTPUT_PATH)
print(f"\n[EXITO] Presentación modificada y guardada en '{OUTPUT_PATH}' con {reemplazados_count} reemplazos realizados.")
